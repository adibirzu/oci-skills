from __future__ import annotations

import json
import re
import uuid
import xml.etree.ElementTree as ET
from copy import copy
from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(".tmp/axian-visuals-update")


def patch_openxml_text(path, replacements):
    tmp = path.with_name(path.stem + "-patched" + path.suffix)
    with ZipFile(path, "r") as src, ZipFile(tmp, "w", compression=ZIP_DEFLATED) as dst:
        for info in src.infolist():
            data = src.read(info.filename)
            if info.filename.endswith(".xml"):
                for old, new in replacements:
                    data = data.replace(old.encode("utf-8"), new.encode("utf-8"))
            dst.writestr(info, data)
    tmp.replace(path)


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def style_snapshot(shape):
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            return {
                "size": run.font.size,
                "name": run.font.name,
                "bold": run.font.bold,
                "italic": run.font.italic,
                "color": None,
            }
    return {"size": None, "name": None, "bold": None, "italic": None, "color": None}


def set_shape_text(shape, text, size=None):
    snap = style_snapshot(shape) if getattr(shape, "has_text_frame", False) else {}
    shape.text = text
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if size:
                run.font.size = Pt(size)
            elif snap.get("size"):
                run.font.size = snap["size"]
            if snap.get("name"):
                run.font.name = snap["name"]
            if snap.get("bold") is not None:
                run.font.bold = snap["bold"]
            if snap.get("italic") is not None:
                run.font.italic = snap["italic"]


def set_named(slide, name, text, size=None):
    matches = [s for s in iter_shapes(slide.shapes) if s.name == name]
    if not matches:
        raise KeyError(f"Shape {name!r} not found")
    set_shape_text(matches[0], text, size=size)


def set_named_color(slide, name, hex_color):
    matches = [s for s in iter_shapes(slide.shapes) if s.name == name]
    if not matches:
        raise KeyError(f"Shape {name!r} not found")
    rgb = RGBColor.from_string(hex_color)
    for paragraph in matches[0].text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = rgb


def replace_exact(slide, old, new, size=None):
    for shape in iter_shapes(slide.shapes):
        if getattr(shape, "has_text_frame", False) and " ".join(shape.text.split()) == " ".join(old.split()):
            set_shape_text(shape, new, size=size)
            return
    raise KeyError(f"Text not found on slide: {old}")


def set_cell(cell, text, size=11):
    cell.text = text
    for p in cell.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(size)


def set_table(slide, rows):
    tables = [s.table for s in slide.shapes if getattr(s, "has_table", False)]
    if len(tables) != 1:
        raise ValueError(f"Expected one table, found {len(tables)}")
    table = tables[0]
    if len(rows) != len(table.rows):
        raise ValueError(f"Expected {len(table.rows)} rows, got {len(rows)}")
    for r, values in zip(table.rows, rows):
        for c, value in zip(r.cells, values):
            set_cell(c, value)


def replace_all(prs, replacements):
    for slide in prs.slides:
        for shape in iter_shapes(slide.shapes):
            if getattr(shape, "has_text_frame", False) and shape.text:
                new = shape.text
                for pattern, repl in replacements:
                    new = re.sub(pattern, repl, new, flags=re.I)
                if new != shape.text:
                    set_shape_text(shape, new)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        new = cell.text
                        for pattern, repl in replacements:
                            new = re.sub(pattern, repl, new, flags=re.I)
                        if new != cell.text:
                            set_cell(cell, new)


def update_safe_harbor(prs):
    slide = prs.slides[1]
    for shape in iter_shapes(slide.shapes):
        if getattr(shape, "has_text_frame", False) and "INSERT MONTH HERE" in shape.text:
            set_shape_text(shape, shape.text.replace("INSERT MONTH HERE", "August"))


def update_management():
    prs = Presentation(ROOT / "management-current.pptx")
    s = prs.slides
    set_named(s[0], "Text Placeholder 2", "Management response | OCI Observability scope for application, API and Kubernetes gaps")
    set_named(s[0], "Text Placeholder 4", "Observability & Security team | RFP AXIAN-IT & Network-OBS-2026 | August 2026")
    for name in ("Text Placeholder 2", "Text Placeholder 3", "Text Placeholder 4"):
        set_named_color(s[0], name, "FFFFFF")
    update_safe_harbor(prs)
    set_named(s[2], "Title 1", "Preserve product ownership. Standardize OCI Observability evidence. Validate every commercial assumption.")

    set_named(s[3], "Title 1", "Three scoped layers, one evidence model")
    set_named(s[3], "Text Placeholder 21", "OCI Observability fills application, API, customer-managed Kubernetes and log-analysis gaps; product teams retain their platforms and remediation.")
    set_named(s[3], "Content Placeholder 2",
              "OCI OBSERVABILITY SERVICES\n• Monitoring metrics, alarms and pipeline health\n• Logging and Connector Hub for approved collection/routing\n• Log Analytics active/archive tiers, parsers, dashboards and LoganAI\n• APM/OpenTelemetry traces, topology, browser and synthetics\n• Notifications for approved HTTPS/email alarm delivery")
    set_named(s[3], "Content Placeholder 3",
              "MANUAL COLLECTION + EXTERNAL BOUNDARIES\n• Customer-managed Kubernetes uses manual Helm/manifests deployment\n• Management Agent/Gateway supports approved hybrid sources\n• Product teams own access, instrumentation, upgrades and remediation\n• External assurance, ITSM and database tooling are out of scope\n• Optional GenAI correlation belongs to the OCI AI Team")

    set_named(s[4], "Title 1", "OCI Observability evidence plane, product-owned action")
    set_named(s[4], "Text Placeholder 2", "Solid flow: telemetry and evidence. Dashed flow: optional advisory correlation. External products and production actions remain product-team owned.")
    box_map = {
        "TextBox 136": "Manual K8s\n+ hybrid collection",
        "TextBox 137": "OCI Observability",
        "TextBox 138": "Sources\nProduct-owned",
        "TextBox 139": "Optional GenAI\nOCI AI Team",
        "TextBox 140": "External ITSM\nOut of scope",
        "TextBox 141": "LoganAI\nAnalyst assistance",
        "TextBox 142": "External workflow\nOut of scope",
        "TextBox 143": "Curated evidence\nDocumented APIs",
        "TextBox 153": "External\nOut of scope",
    }
    for name, text in box_map.items():
        set_named(s[4], name, text, size=10)

    set_named(s[5], "Title 4", "Coverage aligned to the Observability & Security remit")
    set_named(s[5], "Text Placeholder 6", "01 Customer-managed Kubernetes")
    set_named(s[5], "Text Placeholder 5", "Log Analytics Kubernetes Monitoring Solution uses manual Helm/manifests deployment with version, RBAC, egress, resource and rollback validation.")
    set_named(s[5], "Text Placeholder 9", "02 Applications and APIs")
    set_named(s[5], "Text Placeholder 10", "APM and OpenTelemetry add traces, topology, browser visibility and synthetics. Application teams own instrumentation and SLOs.")
    set_named(s[5], "Text Placeholder 15", "03 Logs, metrics and alerting")
    set_named(s[5], "Text Placeholder 16", "Monitoring, Logging, Connector Hub, Log Analytics, Management Agent/Gateway and Notifications provide the shared evidence layer.")

    set_named(s[6], "Title 1", "What remains conditional or requires product-team input")
    set_table(s[6], [
        ["Partial requirement", "Why partial", "Closure requirement"],
        ["Customer-managed K8s deployment", "Distribution, version, RBAC and network path remain open", "Manual Helm/manifests PoV + rollback"],
        ["APM instrumentation", "Application teams must provide service identity and trace propagation", "Priority-app trace-continuity test"],
        ["APM commercial mapping", "300 GB = 1 B95634 active-storage unit is a working presumption", "Oracle quote and observed-volume validation"],
        ["Custom Log Analytics content", "Customer sources may require parsers, labels and dashboards", "Parser accuracy + lifecycle ownership"],
        ["Sizing, resilience & residency", "Volume, retention, regions and recovery remain discovery inputs", "Benchmark + approved decisions"],
    ])

    timeline = {
        "Now Confirm scope & named owners": "Now Confirm sources, product owners and decisions",
        "Month 1 Sources, residency, packs, Senegal pilot": "Month 1 Measure logs, metrics, traces, clusters and residency",
        "Month 1–2 Private foundation & UA core": "Month 1–2 Configure OCI Observability foundation",
        "Month 3 GenAI PoC + guarded automation": "Month 3 Validate LoganAI and optional read-only GenAI PoV",
        "Month 2 OCI O&M + Zurich ServiceNow": "Month 2 Manual Kubernetes monitoring + priority APM",
        "After Phase 1 OPCO onboarding factory": "After Phase 1 Repeatable OPCO onboarding and handover",
    }
    for old, new in timeline.items(): replace_exact(s[7], old, new)

    set_named(s[8], "Title 1", "Decisions required to start the Observability proof of value")
    set_named(s[8], "Text Placeholder 21", "Approve the scoped architecture now; gate pricing and production commitments on measured evidence and the final Oracle quote.")
    set_named(s[8], "Content Placeholder 2",
              "APPROVE NOW\n• OCI Observability-only delivery boundary\n• Manual Kubernetes Monitoring Solution deployment\n• APM/OpenTelemetry for priority applications and APIs\n• Working APM assumption: 300 GB active = 1 × B95634\n• Product teams own access, instrumentation and remediation\n• Optional GenAI correlation stays with the OCI AI Team")
    set_named(s[8], "Content Placeholder 3",
              "GATE BEFORE COMMITMENT\n• Oracle quote confirms B95634 mapping and USD rate\n• Measured log, metric, trace and synthetic volumes\n• Supported Kubernetes distribution/version and private egress\n• Parser coverage, data classes, residency and retention\n• Availability, recovery and monitor-the-monitoring evidence\n• LoganAI/GenAI region, IAM, quality and human-approval controls")

    set_named(s[9], "Title 1", "Approve a scoped Observability proof of value")
    set_named(s[9], "Text Placeholder 3", "Launch a bounded Phase 1: manually onboard representative Kubernetes clusters, instrument priority applications/APIs, validate Log Analytics and APM sizing, and confirm the Oracle commercial mapping before rollout.")

    replace_all(prs, [(r"\bJuly 2026\b", "August 2026")])
    out = ROOT / "AXIAN-Observability-Management-Latest-Answers.pptx"
    prs.save(out)
    return out


def update_engineering():
    prs = Presentation(ROOT / "engineering-current.pptx")
    s = prs.slides
    set_named(s[0], "Text Placeholder 2", "Engineering response | OCI Observability architecture and implementation")
    set_named(s[0], "Text Placeholder 4", "Observability & Security team | Technical architecture | August 2026")
    for name in ("Text Placeholder 2", "Text Placeholder 3", "Text Placeholder 4"):
        set_named_color(s[0], name, "FFFFFF")
    update_safe_harbor(prs)
    set_named(s[2], "Text Placeholder 8", "Collect approved telemetry once; preserve source ownership. Use manual, supported Kubernetes collection. Correlate logs, metrics and traces as governed evidence. Keep production actions product-owned. Confirm pricing, retention, residency and recovery through measured tests.")

    set_named(s[3], "Title 1", "Detailed OCI Observability architecture and communication flows")
    set_named(s[3], "Text Placeholder 2", "Solid = telemetry/evidence. Dashed red = optional GenAI advisory flow. External products and state-changing actions are outside this Observability scope.")
    boxes = {
        "opco-sources": "TELEMETRY SOURCES\nApps • APIs • customer-managed K8s • hybrid hosts",
        "ingestion": "MANUAL / OPEN COLLECTION\nK8s Helm/manifests • OTel • Management Agent/Gateway",
        "unified-assurance": "OCI LOGGING + CONNECTOR HUB\nService/custom logs • governed routing",
        "oci-om": "OCI LOG ANALYTICS\nActive/archive storage • parsers • dashboards",
        "servicenow": "OCI MONITORING\nMetrics • alarms • pipeline health",
        "oem": "OCI APM\nTraces • topology • browser • synthetics",
        "security": "OCI NOTIFICATIONS\nHTTPS/email alarm delivery",
        "loganai": "LOGANAI\nAnalyst assistance • human validated",
        "incident-agent": "OPTIONAL GENAI CORRELATION\nExternal OCI AI Team scope • read-only evidence",
        "product-owners": "PRODUCT TEAMS\nAccess • instrument • remediate • verify",
    }
    for name, text in boxes.items(): set_named(s[3], name, text, size=12)

    lifecycle = {
        "Round Same Side Corner Rectangle 36": "4 Notify and hand off\nNotifications • evidence links • product-owned response and verification",
        "Round Same Side Corner Rectangle 37": "2 Observe and analyze\nMonitoring • Logging • Log Analytics • APM • LoganAI",
        "Round Same Side Corner Rectangle 38": "1 Collect and identify\nManual K8s Helm/manifests • Management Agent/Gateway • OTel • approved logs/metrics",
        "Round Same Side Corner Rectangle 39": "3 Package evidence\nTimestamped logs + metrics + traces + object state + source provenance",
    }
    set_named(s[4], "Title 1", "Telemetry lifecycle: source to evidence to product-owned response")
    for name, text in lifecycle.items(): set_named(s[4], name, text, size=13)

    set_named(s[5], "Title 1", "Customer-managed Kubernetes monitoring is a manual deployment")
    set_named(s[5], "Text Placeholder 21", "Use the documented Log Analytics Kubernetes Monitoring Solution Helm/manifests path; validate each distribution, version and private-network pattern in the PoV.")
    set_named(s[5], "Content Placeholder 2",
              "OCI OBSERVABILITY ROLE\n• Selected Kubernetes logs, metrics and object state\n• Cluster, workload, node and pod views\n• Log Analytics sources, parsers, dashboards and alerts\n• Monitoring metrics and pipeline-health alarms\n• Connector Hub routing\n• Management Agent/Gateway where required")
    set_named(s[5], "Content Placeholder 3",
              "DISCOVERY / VALIDATION BOUNDARY\n• Supported Kubernetes distribution and version\n• Helm/manifests, RBAC and namespace scope\n• Proxy/private egress, DNS and certificates\n• Data classes, filtering and residency\n• CPU/memory/storage footprint\n• Upgrade, rollback, HA and recovery\n• Platform team owns cluster remediation")

    set_named(s[6], "Content Placeholder 3", "OCI LOG ANALYTICS\nCross-source search, clustering, dashboards and retained investigation.\nKUBERNETES SOLUTION\nCustomer-managed Kubernetes logs + metrics + object state with cluster, workload, node and pod views; manual deployment required.")

    set_named(s[7], "Title 1", "APM metering and application-team boundary")
    set_named(s[7], "Text Placeholder 6", "Working commercial presumption: APM active storage is metered through B95634 at one Log Analytics Active Storage Unit per 300 GB; final Oracle quote governs.")
    set_named(s[7], "Content Placeholder 2",
              "OCI APM\n• Application/API traces and service topology\n• Browser and synthetic monitoring\n• OpenTelemetry trace ingestion\n• Approved attributes and sampling\n• Priority journeys and SLO alarms")
    set_named(s[7], "Content Placeholder 3",
              "COMMERCIAL PRESUMPTION\n• 300 GB APM active storage = 1 × B95634\n• First public USD PAYG tier: $372/unit-month\n• API snapshot updated 16 July 2026\n• Confirm part number, conversion and contracted rate\n• Measure active retention before production sizing")
    set_named(s[7], "Content Placeholder 7",
              "OWNERSHIP BOUNDARY\nApplication teams own instrumentation, service IDs, sampling, sensitive-field controls, trace propagation, synthetic journeys and remediation. Observability & Security owns APM domains, standards, dashboards and alert quality.")

    set_named(s[8], "Title 1", "Manual Kubernetes deployment and acceptance checklist")
    set_named(s[8], "Text Placeholder 6", "The Kubernetes Monitoring Solution must be deployed and operated without assuming a managed Kubernetes control plane.")
    set_named(s[8], "Content Placeholder 2",
              "PREREQUISITES\n• Supported distribution/version\n• Cluster RBAC and namespace scope\n• Private egress/proxy/DNS\n• Certificates and agent/gateway placement\n• Approved log/object-state data classes")
    set_named(s[8], "Content Placeholder 3",
              "DEPLOYMENT\n• Manual Helm/manifests path\n• Resource requests/limits\n• Stable cluster/source identity\n• Filtering, parsers and retention\n• Monitoring of collectors, routes and backlog")
    set_named(s[8], "Content Placeholder 7",
              "OPERATE + ACCEPT\n• Upgrade and rollback procedure\n• Agent/gateway HA and reconnect\n• Completeness/latency/loss tests\n• Cluster/workload/node/pod views\n• Platform-team sign-off for access and remediation")

    set_named(s[9], "Content Placeholder 2", "CENTRAL OBSERVABILITY & SECURITY\n• Common telemetry identity and time synchronization\n• Freshness, completeness and parse-success SLOs\n• Least privilege, encryption and access review\n• Active/archive retention boundaries\n• Evidence lineage from source to query/alert\n• Platform audit, cost and capacity controls")
    set_named(s[9], "Content Placeholder 3", "PRODUCT / DATA / COUNTRY OWNERS\n• Approve source onboarding and credentials\n• Classify sensitive and regulated fields\n• Confirm country residency and cross-border rules\n• Instrument applications and maintain service IDs\n• Own thresholds, maintenance and suppression semantics\n• Approve remediation actions and closure evidence")

    set_named(s[10], "Title 1", "LoganAI assists analysis. Optional GenAI correlation remains external.")
    flow11 = {
        "Round Same Side Corner Rectangle 36": "4 Product-owner gate\nReview evidence • diagnose • remediate • verify • rollback",
        "Round Same Side Corner Rectangle 37": "2 LoganAI-assisted investigation\nExplain/summarize logs, clusters and charts • human validation",
        "Round Same Side Corner Rectangle 38": "1 Curated evidence\nMonitoring + Logging + Log Analytics + APM • documented APIs",
        "Round Same Side Corner Rectangle 39": "3 Optional GenAI correlation\nCited timeline + hypotheses + evidence package\nOCI AI Team (Alex Negrea) • read-only by default",
    }
    for name, text in flow11.items(): set_named(s[10], name, text, size=13)

    set_named(s[11], "Title 7", "OCI Observability gaps and closure requirements")
    set_table(s[11], [
        ["Requirement / gap", "Response state", "Closure evidence"],
        ["Exact customer source formats", "Partial — discover versions/schemas", "Interface catalogue + parser compatibility PoV"],
        ["Customer-managed K8s deployment", "Partial — distribution/version/access open", "Manual deployment, upgrade and rollback test"],
        ["APM instrumentation", "Partial — application-team implementation", "Trace continuity, sensitive-field and synthetic tests"],
        ["APM 300 GB / B95634 mapping", "Conditional — commercial presumption", "Final Oracle quote/rate-card confirmation"],
        ["LoganAI and optional GenAI", "Partial — region/IAM/quality validation", "Human evaluation + external AI-team design"],
        ["Volumes, retention and residency", "Partial — measured inputs required", "Benchmark and approved sizing workbook"],
        ["HA, recovery and monitoring health", "Partial — component tests required", "Failure drill, reconnect, replay and restore evidence"],
    ])

    set_named(s[12], "Title 1", "Observability resilience is component-specific, not a single checkbox")
    set_named(s[12], "Text Placeholder 21", "Commit availability and recovery objectives only after collection, routing, storage, APM and alert-delivery failure evidence are agreed.")
    set_named(s[12], "Content Placeholder 2",
              "DESIGN PATTERNS\n• Redundant Management Gateways for private zones\n• Durable configuration and parser lifecycle\n• Log Analytics active/archive retention design\n• Monitoring of collectors, routes, backlog and dropped data\n• APM sampling and agent reconnect behavior\n• Kubernetes collector resource limits and rolling upgrade\n• Product-owned source recovery procedures")
    set_named(s[12], "Content Placeholder 3",
              "ACCEPTANCE EVIDENCE\n• Collector/gateway outage and reconnect\n• Route backlog, replay and duplicate handling\n• Parser failure and data-quality alerting\n• Kubernetes node/namespace collection continuity\n• APM agent outage and trace-gap detection\n• Alert-delivery retry and suppression\n• Restore dashboards/configuration and prove rollback")

    timeline = {
        "Now Owners, sources and decisions": "Now Sources, product owners and decisions",
        "Month 1 Rates, residency, packs, Senegal scope": "Month 1 Measure rates, traces, clusters and residency",
        "Month 1–2 Private foundation and UA core": "Month 1–2 OCI Observability foundation + private collection",
        "Month 3 Agent PoC and guarded automation": "Month 3 LoganAI validation + optional read-only GenAI PoV",
        "Month 2 OCI O&M + Zurich ServiceNow": "Month 2 Manual Kubernetes monitoring + priority APM",
        "Post-pilot OPCO rollout factory and transfer": "Post-pilot OPCO onboarding factory and handover",
    }
    for old, new in timeline.items(): replace_exact(s[13], old, new)

    set_named(s[14], "Text Placeholder 12", "Recommend representative customer-managed Kubernetes clusters, priority applications/APIs and supported hybrid log sources with explicit exit criteria.")
    set_named(s[14], "Content Placeholder 3",
              "1. Ingest agreed logs and metrics with measured freshness, completeness and parse success.\n2. Deploy the Kubernetes Monitoring Solution manually and prove cluster/workload/node/pod visibility.\n3. Demonstrate collector/gateway health, backlog detection, reconnect and rollback.\n4. Instrument priority applications/APIs and prove trace continuity, topology and synthetic journeys.\n5. Validate the 300 GB = 1 × B95634 APM presumption against observed volumes and the final Oracle quote.\n6. Compare human analysis with LoganAI; validate explanations and false conclusions.\n7. If selected, validate the external GenAI prototype with citations, least privilege and human approval.\n8. Produce sizing, cost, residency, resilience and handover evidence for the production decision.")

    set_named(s[15], "Text Placeholder 3", "K8s distribution/version + cluster inventory • priority apps/APIs • telemetry samples • B95634 quote validation • retention/residency • agent/gateway network paths • recovery targets • optional GenAI evaluation", size=14)

    replace_all(prs, [(r"\bJuly 2026\b", "August 2026")])
    out = ROOT / "AXIAN-Observability-Engineering-Latest-Answers.pptx"
    prs.save(out)
    return out


def update_architecture_ppt():
    prs = Presentation(ROOT / "architecture-current.pptx")
    # Rebuild both pages on the Oracle FY26 light-blank layout. The legacy
    # pages used dozens of independent labels which became unreadable when the
    # scoped wording was inserted.
    for old_id in list(prs.slides._sldIdLst):
        old_rid = old_id.rId
        prs.slides._sldIdLst.remove(old_id)
        prs.part.drop_rel(old_rid)
    s1 = prs.slides.add_slide(prs.slide_layouts[10])
    s2 = prs.slides.add_slide(prs.slide_layouts[10])

    def text_box(x, y, w, h, text, size=12, bold=False, color="312D2A", align=PP_ALIGN.LEFT, slide=None):
        slide = slide or s2
        sh = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        sh.text_frame.clear()
        sh.text_frame.word_wrap = True
        p = sh.text_frame.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.name = "Oracle Sans"
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = RGBColor.from_string(color)
        return sh

    def arch_box(x, y, w, h, title, body, fill, stroke, title_color="312D2A", slide=None):
        slide = slide or s2
        sh = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor.from_string(fill)
        sh.line.color.rgb = RGBColor.from_string(stroke)
        sh.line.width = Pt(1.5)
        tf = sh.text_frame
        tf.clear()
        tf.margin_left = Inches(0.14)
        tf.margin_right = Inches(0.14)
        tf.margin_top = Inches(0.12)
        tf.margin_bottom = Inches(0.10)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = title
        r.font.name = "Oracle Sans"
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = RGBColor.from_string(title_color)
        p2 = tf.add_paragraph()
        p2.space_before = Pt(5)
        r2 = p2.add_run()
        r2.text = body
        r2.font.name = "Oracle Sans"
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = RGBColor.from_string("312D2A")
        return sh

    def arrow(x, y, w=0.33, slide=None):
        slide = slide or s2
        sh = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(0.30))
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor.from_string("00758F")
        sh.line.fill.background()
        return sh

    # Page 1: map the broad RFP goals to the bounded Observability response.
    text_box(0.45, 0.42, 12.3, 0.42, "AXIAN requirements and scoped OCI Observability response", 20, True, slide=s1)
    banner = s1.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.45), Inches(1.02), Inches(12.0), Inches(0.52))
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor.from_string("C74634")
    banner.line.fill.background()
    banner.text_frame.clear()
    p = banner.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Group-wide centralized monitoring, observability and operations intelligence"
    r.font.name = "Oracle Sans"
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = RGBColor.from_string("FFFFFF")

    arch_box(0.45, 1.78, 3.66, 1.34, "CENTRALIZED + FEDERATED", "Group and OpCo views • multi-tenancy • scalable onboarding • integration with existing investments", "FFF8F4", "C74634", slide=s1)
    arch_box(4.29, 1.78, 3.66, 1.34, "PROACTIVE OPERATIONS", "Predictive analysis • standardized KPIs • escalation models • evidence-based automation", "FFF6DC", "B36B00", slide=s1)
    arch_box(8.13, 1.78, 4.32, 1.34, "SERVICE + BUSINESS VISIBILITY", "Applications • APIs • logs • metrics • traces • infrastructure and product context", "EAF5F8", "00758F", slide=s1)

    text_box(0.45, 3.43, 12.0, 0.34, "Scoped response from the Observability & Security team", 16, True, color="C74634", slide=s1)
    arch_box(0.45, 3.91, 2.82, 1.55, "OCI APM", "Applications and APIs\nTraces • topology • browser • synthetics\nOpenTelemetry instrumentation", "F2F0FA", "6F4AA8", slide=s1)
    arch_box(3.50, 3.91, 2.82, 1.55, "LOG ANALYTICS + K8S", "Customer-managed Kubernetes\nManual Helm/manifests\nLogs • metrics • object state • LoganAI", "EAF6EF", "2D6A4F", slide=s1)
    arch_box(6.55, 3.91, 2.82, 1.55, "OCI OBSERVABILITY CORE", "Monitoring • Logging • Connector Hub\nManagement Agent/Gateway\nNotifications and evidence links", "EAF5F8", "00758F", slide=s1)
    arch_box(9.60, 3.91, 2.85, 1.55, "OPTIONAL OCI GENAI", "External OCI AI Team workstream\nRead-only evidence correlation\nCitations + human approval", "FFF1EF", "C74634", slide=s1)

    arch_box(0.45, 5.77, 12.0, 0.78, "EXTERNAL PRODUCT BOUNDARY", "Assurance, ITSM, database tooling and shared platform/security foundations remain with their respective product teams and are outside this OCI Observability BOM.", "F3F3F3", "777777", slide=s1)
    text_box(0.45, 6.63, 8.4, 0.25, "No autonomous production change is asserted; product teams own diagnosis, remediation and verification.", 9.5, color="5B5B5B", slide=s1)
    text_box(9.1, 6.63, 3.3, 0.25, "Copyright © 2026, Oracle and/or its affiliates", 8, color="777777", align=PP_ALIGN.RIGHT, slide=s1)

    text_box(0.45, 0.42, 12.3, 0.42, "OCI Observability workflow — customer-managed Kubernetes, applications and logs", 19, True)
    text_box(0.45, 0.95, 12.1, 0.28, "Solid path = telemetry and evidence. Optional GenAI is advisory and external to the OCI Observability BOM.", 10.5, color="5B5B5B")

    arch_box(0.45, 1.52, 2.25, 2.60, "1  TELEMETRY SOURCES", "Applications and APIs\nCustomer-managed Kubernetes\nOCI service and custom logs\nApproved hybrid hosts\n\nProduct-team owned", "FFF8F4", "C74634")
    arch_box(3.04, 1.52, 2.25, 2.60, "2  MANUAL / OPEN COLLECTION", "K8s Helm or manifests\nOpenTelemetry\nManagement Agent / Gateway\nFiltering and source identity\nCollector health and rollback", "EAF5F8", "00758F")
    arch_box(5.63, 1.34, 3.25, 2.96, "3  OCI OBSERVABILITY DATA PLANE", "Monitoring — metrics and alarms\nLogging + Connector Hub — governed routing\nLog Analytics — active/archive, parsers, K8s views and LoganAI\nAPM — traces, topology, browser and synthetics\nNotifications — approved alarm delivery", "EAF6EF", "2D6A4F")
    arch_box(9.22, 1.52, 3.18, 2.60, "4  OPERATIONS EVIDENCE", "Dashboards and alarm context\nCluster, workload, node and pod views\nApplication/API trace topology\nEvidence links and source provenance\nHuman-validated LoganAI assistance", "F2F0FA", "6F4AA8")
    arrow(2.72, 2.62)
    arrow(5.31, 2.62)
    arrow(8.90, 2.62)

    arch_box(0.45, 4.68, 3.68, 1.43, "EXTERNAL PRODUCT BOUNDARY", "Assurance, ITSM, database tooling and shared security foundations remain in their product workstreams and are not included in this BOM.", "F3F3F3", "777777")
    arch_box(4.52, 4.68, 3.48, 1.43, "OPTIONAL GENAI CORRELATION", "External OCI AI Team scope (Alex Negrea). Read-only curated evidence, citations, confidence and human approval before action.", "FFF1EF", "C74634")
    arch_box(8.39, 4.68, 4.01, 1.43, "5  PRODUCT-OWNED ACTION", "Product teams approve access, instrument services, diagnose, remediate, verify and roll back. Observability provides evidence—not autonomous production change.", "FFF6DC", "B36B00")
    arrow(8.02, 5.22)

    text_box(0.45, 6.38, 8.4, 0.30, "Scoped to OCI Observability services for identified application, API, Kubernetes and telemetry gaps.", 9.5, color="5B5B5B")
    text_box(9.1, 6.38, 3.3, 0.30, "Copyright © 2026, Oracle and/or its affiliates", 8, color="777777", align=PP_ALIGN.RIGHT)

    out = ROOT / "Axian-Telecom-SolutionArch-31July2026-Latest-Answers.pptx"
    prs.save(out)
    patch_openxml_text(out, [("Confidential - Oracle Restricted \\Employees Only", "Confidential - Oracle Restricted | Employees Only")])
    return out


def mx_cell(parent, cell_id, value, x, y, w, h, style, parent_id="1"):
    cell = ET.SubElement(parent, "mxCell", {"id": cell_id, "value": value, "style": style, "vertex": "1", "parent": parent_id})
    ET.SubElement(cell, "mxGeometry", {"x": str(x), "y": str(y), "width": str(w), "height": str(h), "as": "geometry"})
    return cell


def mx_edge(parent, cell_id, source, target, label, dashed=False):
    style = "edgeStyle=orthogonalEdgeStyle;rounded=0;orthogonalLoop=1;jettySize=auto;html=1;endArrow=block;endFill=1;strokeWidth=2;strokeColor=#00688B;"
    if dashed:
        style += "dashed=1;strokeColor=#C74634;"
    cell = ET.SubElement(parent, "mxCell", {"id": cell_id, "value": label, "style": style, "edge": "1", "parent": "1", "source": source, "target": target})
    ET.SubElement(cell, "mxGeometry", {"relative": "1", "as": "geometry"})


def add_page(mxfile, page_id, name, title, nodes, edges):
    diagram = ET.SubElement(mxfile, "diagram", {"id": page_id, "name": name})
    model = ET.SubElement(diagram, "mxGraphModel", {"dx": "1600", "dy": "1000", "grid": "1", "gridSize": "10", "page": "1", "pageWidth": "1600", "pageHeight": "1000", "math": "0", "shadow": "0"})
    root = ET.SubElement(model, "root")
    ET.SubElement(root, "mxCell", {"id": "0"})
    ET.SubElement(root, "mxCell", {"id": "1", "parent": "0"})
    mx_cell(root, "title", title, 60, 25, 1480, 55, "text;html=1;strokeColor=none;fillColor=none;fontSize=26;fontStyle=1;fontColor=#312D2A;align=left;verticalAlign=middle;")
    for node in nodes:
        mx_cell(root, *node)
    for edge in edges:
        mx_edge(root, *edge)


def build_drawio():
    mxfile = ET.Element("mxfile", {"host": "app.diagrams.net", "modified": "2026-08-03T00:00:00.000Z", "agent": "Codex", "version": "24.7.17", "type": "device"})
    box = "rounded=1;whiteSpace=wrap;html=1;strokeWidth=2;strokeColor=#00688B;fillColor=#EAF5F8;fontColor=#312D2A;fontSize=16;align=center;verticalAlign=middle;spacing=8;"
    redbox = "rounded=1;whiteSpace=wrap;html=1;strokeWidth=2;strokeColor=#C74634;fillColor=#FFF1EF;fontColor=#312D2A;fontSize=16;align=center;verticalAlign=middle;spacing=8;"
    greenbox = "rounded=1;whiteSpace=wrap;html=1;strokeWidth=2;strokeColor=#2D6A4F;fillColor=#EAF6EF;fontColor=#312D2A;fontSize=16;align=center;verticalAlign=middle;spacing=8;"
    graybox = "rounded=1;whiteSpace=wrap;html=1;strokeWidth=2;strokeColor=#777777;fillColor=#F2F2F2;fontColor=#312D2A;fontSize=16;align=center;verticalAlign=middle;spacing=8;"

    add_page(mxfile, "obs-e2e", "01 - OCI Observability Architecture", "AXIAN — OCI Observability Evidence Architecture (Scoped)", [
        ("src", "<b>TELEMETRY SOURCES</b><br>Applications • APIs<br>Customer-managed Kubernetes<br>OCI service/custom logs • hybrid hosts<br><i>Product-team owned</i>", 70, 180, 260, 260, box),
        ("collect", "<b>MANUAL / OPEN COLLECTION</b><br>K8s Helm/manifests<br>OpenTelemetry<br>Management Agent/Gateway<br>Filtering • source identity • health", 390, 180, 270, 260, box),
        ("plane", "<b>OCI OBSERVABILITY DATA PLANE</b><br>Monitoring<br>Logging + Connector Hub<br>Log Analytics active/archive<br>APM + synthetics<br>Notifications", 730, 150, 310, 320, greenbox),
        ("outcome", "<b>OPERATOR OUTCOMES</b><br>Dashboards • alarms • evidence links<br>K8s cluster/workload/node/pod views<br>Trace topology • journey health<br>LoganAI-assisted investigation", 1110, 180, 300, 260, box),
        ("genai", "<b>OPTIONAL GENAI CORRELATION</b><br>External OCI AI Team scope<br>Read-only curated evidence<br>Citations • confidence • human review", 760, 570, 280, 180, redbox),
        ("owners", "<b>PRODUCT TEAMS</b><br>Approve access • instrument<br>Diagnose • remediate • verify • rollback", 1130, 570, 270, 180, graybox),
        ("boundary", "<b>EXTERNAL PRODUCT BOUNDARY</b><br>Assurance • ITSM • database tooling<br>Networking/security foundation<br>Not included in this BOM", 70, 570, 590, 180, graybox),
    ], [
        ("e1", "src", "collect", "logs • metrics • traces • object state", False),
        ("e2", "collect", "plane", "supported ingest / documented APIs", False),
        ("e3", "plane", "outcome", "queries • dashboards • alarms", False),
        ("e4", "plane", "genai", "curated evidence only", True),
        ("e5", "genai", "owners", "advisory package", True),
        ("e6", "outcome", "owners", "evidence + notification", False),
    ])

    add_page(mxfile, "k8s-manual", "02 - Manual Kubernetes Monitoring", "Customer-Managed Kubernetes — Manual Log Analytics Monitoring Deployment", [
        ("k8s", "<b>KUBERNETES CLUSTER</b><br>Supported distribution/version<br>Namespaces • nodes • workloads<br>Platform-team ownership", 80, 180, 270, 240, box),
        ("deploy", "<b>MANUAL DEPLOYMENT</b><br>Documented Helm/manifests<br>RBAC • namespace scope<br>Proxy/private egress • certificates<br>Resource requests/limits", 430, 150, 300, 300, redbox),
        ("collect2", "<b>COLLECTION</b><br>Selected logs + metrics<br>Kubernetes object state<br>Stable cluster/source identity<br>Management Agent/Gateway where needed", 810, 150, 300, 300, box),
        ("la", "<b>OCI LOG ANALYTICS</b><br>Kubernetes Monitoring Solution<br>Cluster • workload • node • pod views<br>Parsers • dashboards • alerts", 1190, 180, 290, 240, greenbox),
        ("accept", "<b>ACCEPTANCE</b><br>Completeness • latency • loss<br>Collector/gateway health<br>Upgrade • reconnect • rollback<br>Data residency and retention", 620, 570, 420, 190, graybox),
    ], [
        ("k1", "k8s", "deploy", "cluster access", False),
        ("k2", "deploy", "collect2", "install + configure", False),
        ("k3", "collect2", "la", "logs • metrics • object state", False),
        ("k4", "la", "accept", "evidence", False),
        ("k5", "accept", "deploy", "validated lifecycle", True),
    ])

    add_page(mxfile, "apm", "03 - APM and Commercial Assumption", "OCI APM — Telemetry Flow and 300 GB Active-Storage Presumption", [
        ("apps", "<b>PRIORITY APPS + APIs</b><br>Java • .NET • PHP • middleware<br>Browser and synthetic journeys<br>Application-team ownership", 100, 190, 280, 240, box),
        ("otel", "<b>INSTRUMENTATION</b><br>APM agents + OpenTelemetry<br>Service/resource attributes<br>Trace propagation • sampling<br>Sensitive-field controls", 470, 170, 290, 280, box),
        ("apm2", "<b>OCI APM</b><br>Traces • service maps<br>Browser • synthetics • topology<br>SLO indicators and alarms", 850, 190, 280, 240, greenbox),
        ("price", "<b>COMMERCIAL WORKING PRESUMPTION</b><br>300 GB APM active storage = 1 × B95634<br>First public USD PAYG tier: $372/unit-month<br>Final Oracle quote/rate card governs", 1210, 160, 310, 300, redbox),
        ("measure", "<b>SIZING EVIDENCE</b><br>Spans • sessions • synthetics<br>Sampling • active retention<br>Non-production/DR • regions<br>Measured before production commitment", 600, 590, 420, 190, graybox),
    ], [
        ("a1", "apps", "otel", "instrument", False),
        ("a2", "otel", "apm2", "OTLP/APM telemetry", False),
        ("a3", "apm2", "price", "active-storage consumption", False),
        ("a4", "apm2", "measure", "observed volumes", False),
        ("a5", "measure", "price", "quote validation", True),
    ])

    add_page(mxfile, "genai", "04 - Optional GenAI Boundary", "Optional GenAI Correlation — External AI-Team Dependency", [
        ("evidence", "<b>CURATED OCI OBSERVABILITY EVIDENCE</b><br>Monitoring • Logging • Log Analytics • APM<br>Source IDs • timestamps • evidence links", 100, 210, 340, 240, greenbox),
        ("api", "<b>DOCUMENTED / ALLOWLISTED APIS</b><br>Least privilege • read-only by default<br>Redaction • timeout • audit • provenance", 520, 210, 330, 240, box),
        ("ai", "<b>OPTIONAL GENAI COORDINATOR</b><br>OCI AI Team (Alex Negrea)<br>Cited timeline • hypotheses • confidence<br>Evaluation and human review required", 930, 180, 350, 300, redbox),
        ("owner", "<b>PRODUCT OWNER</b><br>Accept/reject recommendation<br>Diagnose • remediate • verify • rollback", 1360, 210, 230, 240, graybox),
        ("logan", "<b>LOGANAI</b><br>Analyst assistance inside Log Analytics<br>Not an agentic incident commander", 620, 590, 360, 160, box),
    ], [
        ("g1", "evidence", "api", "approved evidence", False),
        ("g2", "api", "ai", "bounded retrieval", True),
        ("g3", "ai", "owner", "advisory package", True),
        ("g4", "evidence", "logan", "logs/metrics for analysis", False),
        ("g5", "logan", "ai", "reviewed evidence only", True),
    ])

    tree = ET.ElementTree(mxfile)
    ET.indent(tree, space="  ")
    out = ROOT / "AXIAN-Observability-Architectures-Latest-Answers.drawio"
    tree.write(out, encoding="utf-8", xml_declaration=True)
    return out


def excal_id(prefix):
    return f"{prefix}-{uuid.uuid4().hex[:10]}"


def ex_rect(elements, x, y, w, h, fill, stroke="#00688B"):
    rid = excal_id("rect")
    elements.append({"id": rid, "type": "rectangle", "x": x, "y": y, "width": w, "height": h, "angle": 0, "strokeColor": stroke, "backgroundColor": fill, "fillStyle": "solid", "strokeWidth": 2, "strokeStyle": "solid", "roughness": 1, "opacity": 100, "groupIds": [], "frameId": None, "roundness": {"type": 3}, "seed": 1, "version": 1, "versionNonce": 1, "isDeleted": False, "boundElements": [], "updated": 1, "link": None, "locked": False})
    return rid


def ex_text(elements, x, y, w, h, text, size=22, color="#312D2A", align="center"):
    tid = excal_id("text")
    elements.append({"id": tid, "type": "text", "x": x, "y": y, "width": w, "height": h, "angle": 0, "strokeColor": color, "backgroundColor": "transparent", "fillStyle": "solid", "strokeWidth": 1, "strokeStyle": "solid", "roughness": 1, "opacity": 100, "groupIds": [], "frameId": None, "roundness": None, "seed": 1, "version": 1, "versionNonce": 1, "isDeleted": False, "boundElements": [], "updated": 1, "link": None, "locked": False, "text": text, "fontSize": size, "fontFamily": 5, "textAlign": align, "verticalAlign": "middle", "containerId": None, "originalText": text, "autoResize": False, "lineHeight": 1.25})
    return tid


def ex_arrow(elements, x1, y1, x2, y2, dashed=False, color="#00688B"):
    aid = excal_id("arrow")
    elements.append({"id": aid, "type": "arrow", "x": x1, "y": y1, "width": x2-x1, "height": y2-y1, "angle": 0, "strokeColor": color, "backgroundColor": "transparent", "fillStyle": "solid", "strokeWidth": 3, "strokeStyle": "dashed" if dashed else "solid", "roughness": 1, "opacity": 100, "groupIds": [], "frameId": None, "roundness": {"type": 2}, "seed": 1, "version": 1, "versionNonce": 1, "isDeleted": False, "boundElements": [], "updated": 1, "link": None, "locked": False, "points": [[0,0],[x2-x1,y2-y1]], "lastCommittedPoint": None, "startBinding": None, "endBinding": None, "startArrowhead": None, "endArrowhead": "arrow", "elbowed": False})
    return aid


def build_excalidraw():
    elements = []
    ex_text(elements, 70, 30, 1500, 70, "AXIAN OCI Observability — Scoped Components and Communication Flows", 34, align="left")
    nodes = [
        (80, 180, 260, 250, "#EAF5F8", "#00688B", "TELEMETRY SOURCES\n\nApps + APIs\nCustomer-managed Kubernetes\nOCI/custom logs + hybrid hosts\n\nProduct-team owned"),
        (420, 180, 280, 250, "#EAF5F8", "#00688B", "MANUAL / OPEN COLLECTION\n\nK8s Helm/manifests\nOpenTelemetry\nManagement Agent/Gateway\nFiltering + source health"),
        (780, 150, 340, 310, "#EAF6EF", "#2D6A4F", "OCI OBSERVABILITY DATA PLANE\n\nMonitoring\nLogging + Connector Hub\nLog Analytics active/archive\nAPM + synthetics\nNotifications"),
        (1200, 180, 300, 250, "#EAF5F8", "#00688B", "OPERATOR OUTCOMES\n\nDashboards + alarms\nK8s cluster/workload/node/pod views\nTrace topology + journey health\nLoganAI-assisted investigation"),
        (795, 570, 310, 210, "#FFF1EF", "#C74634", "OPTIONAL GENAI CORRELATION\n\nExternal OCI AI Team scope\nRead-only curated evidence\nCitations + confidence\nHuman review required"),
        (1210, 570, 290, 210, "#F2F2F2", "#777777", "PRODUCT TEAMS\n\nApprove access + instrument\nDiagnose + remediate\nVerify + rollback"),
        (80, 570, 620, 210, "#F2F2F2", "#777777", "EXTERNAL PRODUCT BOUNDARY\n\nAssurance + ITSM + database tooling\nNetworking/security foundation\nOutside the OCI Observability BOM"),
    ]
    for x,y,w,h,fill,stroke,text in nodes:
        ex_rect(elements,x,y,w,h,fill,stroke)
        ex_text(elements,x+15,y+15,w-30,h-30,text,20)
    ex_arrow(elements,340,305,420,305)
    ex_arrow(elements,700,305,780,305)
    ex_arrow(elements,1120,305,1200,305)
    ex_arrow(elements,950,460,950,570,True,"#C74634")
    ex_arrow(elements,1105,675,1210,675,True,"#C74634")
    ex_arrow(elements,1350,430,1350,570)
    ex_text(elements,80,830,1450,80,"Solid arrows = telemetry/evidence. Dashed red arrows = optional advisory correlation. Production access and remediation remain product-team owned.",18,"#5B5B5B","left")
    out = ROOT / "AXIAN-Observability-Architecture-Latest-Answers.excalidraw"
    out.write_text(json.dumps({"type":"excalidraw","version":2,"source":"https://excalidraw.com","elements":elements,"appState":{"viewBackgroundColor":"#ffffff","gridSize":20},"files":{}}, indent=2))
    return out


if __name__ == "__main__":
    outputs = [update_management(), update_engineering(), update_architecture_ppt(), build_drawio(), build_excalidraw()]
    for output in outputs:
        print(output)
